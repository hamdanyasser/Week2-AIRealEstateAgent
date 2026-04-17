"""Generate the AI Real Estate Agent presentation deck.

Run:    python make_deck.py
Output: AI_Real_Estate_Agent.pptx  (9 slides, 16:9, with fade transitions)

A character-driven product brief. Meet Layla, a first-time home buyer.
The deck walks through her story in nine scenes — problem, idea,
transformation, the helpers behind the scenes, trust, demo, and vision.

Fonts: Georgia, Segoe UI, Consolas — all Windows 11 defaults.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ============================================================
# DESIGN TOKENS
# ============================================================
CANVAS     = RGBColor(0xF6, 0xF1, 0xE8)
PAPER      = RGBColor(0xFB, 0xF7, 0xEF)
SURFACE    = RGBColor(0xEF, 0xE7, 0xDA)
DEEP       = RGBColor(0xE6, 0xDA, 0xC5)
GHOST      = RGBColor(0xDF, 0xD4, 0xBE)
SAND       = RGBColor(0xEA, 0xDF, 0xC9)

INK        = RGBColor(0x2F, 0x3A, 0x39)
INK_SOFT   = RGBColor(0x4A, 0x55, 0x53)
MUTED      = RGBColor(0x6B, 0x75, 0x6F)
FAINT      = RGBColor(0x9B, 0xA1, 0x9A)

TERRACOTTA = RGBColor(0xC8, 0x55, 0x3D)
TERRA_DEEP = RGBColor(0x9E, 0x3F, 0x2A)
TERRA_SOFT = RGBColor(0xE8, 0xA8, 0x97)
TEAL       = RGBColor(0x2F, 0x6F, 0x6A)
TEAL_SOFT  = RGBColor(0x8E, 0xB5, 0xB1)
BRASS      = RGBColor(0xB8, 0x8B, 0x4A)
BRASS_SOFT = RGBColor(0xDC, 0xC3, 0x92)

HAIRLINE   = RGBColor(0xCF, 0xC5, 0xB4)

SERIF = "Georgia"
SANS  = "Segoe UI"
MONO  = "Consolas"

TOTAL_SCENES = 9

# ============================================================
# DECK
# ============================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

P_NS  = "http://schemas.openxmlformats.org/presentationml/2006/main"

# ============================================================
# PRIMITIVES
# ============================================================
def paint_canvas(slide, color=CANVAS):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    tree = r._element.getparent()
    tree.remove(r._element); tree.insert(2, r._element)

def new_slide(bg=CANVAS):
    s = prs.slides.add_slide(BLANK)
    paint_canvas(s, bg)
    add_fade_transition(s)
    return s

def _frame(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf

def text(slide, l, t, w, h, s,
         font=SANS, size=18, bold=False, italic=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb, tf = _frame(slide, l, t, w, h)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = s
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def multi(slide, l, t, w, h, lines,
          font=SANS, size=18, bold=False, italic=False, color=INK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25,
          space_after=0):
    tb, tf = _frame(slide, l, t, w, h)
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return tb

def hairline(slide, x1, y1, x2, y2, color=HAIRLINE, w=0.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(w)
    return c

def rect(slide, l, t, w, h, fill=None, stroke=None, stroke_w=0.75):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if stroke is None: r.line.fill.background()
    else: r.line.color.rgb = stroke; r.line.width = Pt(stroke_w)
    return r

def rrect(slide, l, t, w, h, fill=None, stroke=None, stroke_w=0.75, radius=0.1):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    r.adjustments[0] = radius
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if stroke is None: r.line.fill.background()
    else: r.line.color.rgb = stroke; r.line.width = Pt(stroke_w)
    return r

def oval(slide, l, t, w, h, fill=None, stroke=None, stroke_w=0.75):
    r = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if stroke is None: r.line.fill.background()
    else: r.line.color.rgb = stroke; r.line.width = Pt(stroke_w)
    return r

def triangle(slide, l, t, w, h, fill=None, stroke=None):
    r = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, l, t, w, h)
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if stroke is None: r.line.fill.background()
    else: r.line.color.rgb = stroke
    return r

def right_arrow(slide, l, t, w, h, fill=TERRACOTTA, stroke=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if stroke is None: r.line.fill.background()
    else: r.line.color.rgb = stroke
    return r

# ============================================================
# CUSTOM ILLUSTRATIONS (built from shapes)
# ============================================================
def person(slide, cx, cy, h, color=TERRACOTTA):
    """Simple centred person silhouette — head + shoulders."""
    head_d = int(h * 0.38)
    head_y = int(cy - h/2)
    oval(slide, cx - head_d/2, head_y, head_d, head_d, fill=color)

    body_w = int(h * 0.72)
    body_h = int(h * 0.52)
    body_y = head_y + head_d + int(h * 0.04)
    rrect(slide, cx - body_w/2, body_y, body_w, body_h,
          fill=color, radius=0.45)

def house(slide, cx, cy, w, h, color=INK, door_color=CANVAS, window_color=None):
    """Simple house — triangle roof, square body, door, two windows."""
    roof_h = int(h * 0.44)
    triangle(slide, cx - w/2, cy - h/2, w, roof_h, fill=color)

    body_w = int(w * 0.82)
    body_x = cx - body_w/2
    body_y = cy - h/2 + int(roof_h * 0.92)
    body_h = int((cy + h/2) - body_y)
    rect(slide, body_x, body_y, body_w, body_h, fill=color)

    # door
    door_w = int(body_w * 0.20)
    door_h = int(body_h * 0.45)
    rect(slide, cx - door_w/2, cy + h/2 - door_h,
         door_w, door_h, fill=door_color)

    # two small windows flanking the door
    if window_color is None:
        window_color = door_color
    win_w = int(body_w * 0.18)
    win_h = int(body_h * 0.22)
    win_y = body_y + int(body_h * 0.18)
    rect(slide, body_x + int(body_w * 0.12), win_y, win_w, win_h, fill=window_color)
    rect(slide, body_x + body_w - int(body_w * 0.12) - win_w, win_y, win_w, win_h, fill=window_color)

def speech_bubble(slide, x, y, w, h, fill=PAPER, stroke=HAIRLINE, stroke_w=1.0):
    """Rounded-rectangular callout. Draws a clean bubble with a small tail."""
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT, x, y, w, h)
    r.adjustments[0] = -0.15    # tail x (-ve = left of shape)
    r.adjustments[1] =  0.55    # tail y (0..1 along shape height)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.color.rgb = stroke
    r.line.width = Pt(stroke_w)
    return r

def target_icon(slide, cx, cy, d, color=TERRACOTTA):
    """Concentric bullseye — three rings."""
    for i, frac in enumerate([1.0, 0.66, 0.33]):
        ring_d = int(d * frac)
        if i == 2:
            oval(slide, cx - ring_d/2, cy - ring_d/2, ring_d, ring_d, fill=color)
        else:
            oval(slide, cx - ring_d/2, cy - ring_d/2, ring_d, ring_d,
                 fill=PAPER, stroke=color, stroke_w=1.5)

def house_stack_icon(slide, cx, cy, w, h, color=TEAL):
    """Three small houses arranged in a neat row — 'many homes'."""
    hw = int(w * 0.3)
    hh = int(h * 0.8)
    gap = (w - 3 * hw) / 2
    for i in range(3):
        hx = cx - w/2 + i * (hw + gap) + hw/2
        house(slide, hx, cy, hw, hh, color=color, door_color=PAPER, window_color=PAPER)

def chat_icon(slide, cx, cy, w, h, color=BRASS):
    """Small speech bubble with a check — 'understood'."""
    bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
                                    cx - w/2, cy - h/2, w, h)
    bubble.adjustments[0] = -0.2
    bubble.adjustments[1] =  0.9
    bubble.fill.solid(); bubble.fill.fore_color.rgb = color
    bubble.line.fill.background()
    # check mark
    text(slide, cx - w/2, cy - h/2 - Inches(0.05), w, h,
         "\u2713", font=SANS, size=40, bold=True, color=PAPER,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# TRANSITIONS
# ============================================================
def add_fade_transition(slide, speed="med"):
    """Add a subtle fade transition to this slide."""
    slide_el = slide._element
    for t in slide_el.findall("{%s}transition" % P_NS):
        slide_el.remove(t)
    transition = etree.SubElement(slide_el, "{%s}transition" % P_NS)
    transition.set("spd", speed)
    etree.SubElement(transition, "{%s}fade" % P_NS)

# ============================================================
# CHROME — appears on every inner slide
# ============================================================
def scene_dots(slide, current):
    """9-dot scene-progress indicator across the top-centre."""
    total = TOTAL_SCENES
    dot = Inches(0.11)
    gap = Inches(0.22)
    step = dot + gap
    total_w = step * (total - 1) + dot
    start_x = (SW - total_w) / 2
    y = Inches(0.51)
    for i in range(total):
        x = start_x + step * i
        if i == current:
            oval(slide, x, y, dot, dot, fill=TERRACOTTA)
        elif i < current:
            oval(slide, x, y, dot, dot, fill=PAPER, stroke=TERRACOTTA, stroke_w=1.0)
        else:
            oval(slide, x, y, dot, dot, fill=HAIRLINE)

def chrome(slide, scene, section):
    # masthead top-left
    rect(slide, Inches(0.6), Inches(0.48), Inches(0.12), Inches(0.12), fill=TERRACOTTA)
    text(slide, Inches(0.8), Inches(0.45), Inches(5), Inches(0.22),
         "AIRE  \u00b7  A PRODUCT BRIEF",
         font=MONO, size=9, color=INK)

    # centred scene dots
    scene_dots(slide, scene)

    # edition top-right
    text(slide, Inches(8.5), Inches(0.45), Inches(4.23), Inches(0.22),
         f"SCENE {scene+1:02d} / {TOTAL_SCENES:02d}  \u00b7  {section.upper()}",
         font=MONO, size=9, color=TERRACOTTA, bold=True, align=PP_ALIGN.RIGHT)

    hairline(slide, Inches(0.6), Inches(0.78), Inches(12.73), Inches(0.78))

    # bottom hairline + footer
    hairline(slide, Inches(0.6), Inches(7.02), Inches(12.73), Inches(7.02))
    text(slide, Inches(0.6), Inches(7.12), Inches(6), Inches(0.22),
         "AIRE  \u00b7  AI REAL ESTATE AGENT",
         font=MONO, size=8, color=FAINT)
    text(slide, Inches(6.5), Inches(7.12), Inches(6.23), Inches(0.22),
         "SE FACTORY  \u00b7  AIE BOOTCAMP  \u00b7  WEEK 2  \u00b7  2026",
         font=MONO, size=8, color=FAINT, align=PP_ALIGN.RIGHT)


# ============================================================
# SCENE 1 — COVER
# ============================================================
def slide_cover():
    s = new_slide()

    # --- masthead (cover-specific) ---
    rect(s, Inches(0.6), Inches(0.48), Inches(0.14), Inches(0.14), fill=TERRACOTTA)
    text(s, Inches(0.84), Inches(0.44), Inches(6), Inches(0.25),
         "AIRE  \u00b7  AI REAL ESTATE AGENT",
         font=MONO, size=9, color=INK, bold=True)

    # scene dots (scene 0 = cover)
    scene_dots(s, 0)

    text(s, Inches(8.5), Inches(0.44), Inches(4.23), Inches(0.25),
         "SPRING 2026  \u00b7  EDITION 01",
         font=MONO, size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    hairline(s, Inches(0.6), Inches(0.8), Inches(12.73), Inches(0.8))

    # --- kicker ---
    rect(s, Inches(0.6), Inches(1.55), Inches(0.06), Inches(0.85), fill=TERRACOTTA)
    text(s, Inches(0.88), Inches(1.55), Inches(10), Inches(0.3),
         "A STORY IN NINE SCENES",
         font=MONO, size=12, color=TERRACOTTA, bold=True)
    text(s, Inches(0.88), Inches(1.98), Inches(11), Inches(0.3),
         "About why pricing a home shouldn\u2019t require a spreadsheet.",
         font=SERIF, size=15, italic=True, color=INK_SOFT)

    # --- giant headline ---
    multi(s, Inches(0.6), Inches(2.95), Inches(12.2), Inches(3.2),
          ["Describe a home.",
           "Get a price."],
          font=SERIF, size=108, color=INK, line_spacing=0.98)

    # --- decorative composition: person looking at a house ---
    person(s, Inches(10.4), Inches(6.3), Inches(0.9), color=TERRACOTTA)
    house(s, Inches(11.9),  Inches(6.3), Inches(1.4), Inches(1.1),
          color=INK, door_color=CANVAS, window_color=TERRA_SOFT)

    # a small dotted gaze line between person and house
    for i in range(6):
        dx = Inches(10.9) + Inches(0.15) * i
        oval(s, dx, Inches(6.28), Inches(0.05), Inches(0.05), fill=MUTED)

    # --- bottom meta ---
    hairline(s, Inches(0.6), Inches(6.8), Inches(9.5), Inches(6.8))
    text(s, Inches(0.6), Inches(6.9), Inches(9), Inches(0.25),
         "A three-minute product brief, told through one buyer\u2019s story.",
         font=SERIF, size=13, italic=True, color=INK_SOFT)
    text(s, Inches(0.6), Inches(7.2), Inches(9), Inches(0.25),
         "SE FACTORY  \u00b7  AIE BOOTCAMP  \u00b7  WEEK 2",
         font=MONO, size=9, color=MUTED)


# ============================================================
# SCENE 2 — MEET LAYLA
# ============================================================
def slide_meet_sarah():
    s = new_slide()
    chrome(s, 1, "Meet Layla")

    # kicker
    text(s, Inches(0.6), Inches(1.2), Inches(10), Inches(0.25),
         "THE PERSON AT THE CENTRE OF THIS STORY",
         font=MONO, size=9, color=MUTED, bold=True)

    # headline
    multi(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.5),
          ["Meet Layla."],
          font=SERIF, size=68, color=INK, line_spacing=1.0)

    # --- LEFT column: the character ---
    # background circle (warm blush)
    oval(s, Inches(0.8), Inches(3.4), Inches(3.6), Inches(3.6),
         fill=TERRA_SOFT)
    # person silhouette
    person(s, Inches(2.6), Inches(5.2), Inches(2.6), color=TERRACOTTA)
    # name label
    text(s, Inches(0.6), Inches(7.1 - 0.35), Inches(4), Inches(0.3),
         "Layla, 32", font=MONO, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # --- RIGHT column: her speech bubble ---
    bubble_x = Inches(5.3)
    bubble_y = Inches(3.6)
    bubble_w = Inches(7.3)
    bubble_h = Inches(2.9)
    speech_bubble(s, bubble_x, bubble_y, bubble_w, bubble_h,
                  fill=PAPER, stroke=HAIRLINE, stroke_w=1.25)

    # big opening quote mark inside the bubble
    text(s, bubble_x + Inches(0.3), bubble_y + Inches(-0.2),
         Inches(1.0), Inches(1.2),
         "\u201C", font=SERIF, size=80, color=TERRACOTTA)

    # the quote itself
    multi(s, bubble_x + Inches(0.65), bubble_y + Inches(0.55),
          bubble_w - Inches(1.0), bubble_h - Inches(0.9),
          ["I just want to know",
           "how much a house is worth.",
           "Why does it need a spreadsheet?"],
          font=SERIF, size=24, italic=True, color=INK, line_spacing=1.3)

    # attribution
    text(s, bubble_x + Inches(0.65), bubble_y + bubble_h + Inches(0.2),
         bubble_w - Inches(1.0), Inches(0.3),
         "\u2014  Layla, first-time home buyer",
         font=SERIF, size=13, italic=True, color=MUTED)


# ============================================================
# SCENE 3 — THE WALL
# ============================================================
def slide_the_wall():
    s = new_slide()
    chrome(s, 2, "The wall")

    # kicker
    text(s, Inches(0.6), Inches(1.2), Inches(10), Inches(0.25),
         "WHAT HAPPENS WHEN SHE OPENS A PRICING TOOL",
         font=MONO, size=9, color=MUTED, bold=True)

    # headline (left column)
    multi(s, Inches(0.6), Inches(1.6), Inches(7), Inches(2.0),
          ["She has to fill", "this out first."],
          font=SERIF, size=50, color=INK, line_spacing=1.03)

    # narrator line
    hairline(s, Inches(0.6), Inches(4.0), Inches(7), Inches(4.0),
             color=TERRACOTTA, w=1.25)
    multi(s, Inches(0.6), Inches(4.15), Inches(7), Inches(1.5),
          ["80 fields.",
           "Most she\u2019s never heard of.",
           "Half need a professional to answer."],
          font=SERIF, size=18, italic=True, color=INK_SOFT, line_spacing=1.45)

    # pull quote
    text(s, Inches(0.6), Inches(5.95), Inches(7), Inches(0.3),
         "SHE GIVES UP.",
         font=MONO, size=11, color=TERRACOTTA, bold=True)
    text(s, Inches(0.6), Inches(6.3), Inches(7), Inches(0.4),
         "This is the moment real estate lost her.",
         font=SERIF, size=16, italic=True, color=INK)

    # --- RIGHT column: the impossible form ---
    col_x = Inches(7.9)
    col_w = Inches(4.83)

    text(s, col_x, Inches(1.2), col_w, Inches(0.25),
         "PROPERTY VALUATION FORM", font=MONO, size=9, color=MUTED, bold=True,
         align=PP_ALIGN.CENTER)
    hairline(s, col_x, Inches(1.52), col_x + col_w, Inches(1.52),
             color=HAIRLINE, w=1.0)

    # technical field names — real Ames columns, on purpose overwhelming
    fields = [
        "MSSubClass",      "MSZoning",        "LotFrontage",
        "LotArea",         "OverallQual",     "OverallCond",
        "YearBuilt",       "MasVnrType",      "BsmtFinType1",
        "HeatingQC",       "CentralAir",      "1stFlrSF",
        "GrLivArea",       "GarageCars",      "PavedDrive",
    ]
    row_h = Inches(0.27)
    y0 = Inches(1.7)
    for i, name in enumerate(fields):
        # label
        text(s, col_x, y0 + row_h * i, Inches(1.85), Inches(0.22),
             name, font=MONO, size=8, color=MUTED)
        # field box (fade rows get fainter)
        fade = 1.0 - (i / len(fields)) * 0.3
        box_fill = SURFACE
        rect(s, col_x + Inches(1.9), y0 + row_h * i + Inches(0.02),
             Inches(2.9), Inches(0.18),
             fill=box_fill, stroke=HAIRLINE, stroke_w=0.5)
        # dash inside
        text(s, col_x + Inches(2.0), y0 + row_h * i - Inches(0.02),
             Inches(2.7), Inches(0.2),
             "\u2014", font=MONO, size=8, color=FAINT)

    # fade-to-nothing gradient: 3 progressively fainter rows
    y_fade = y0 + row_h * len(fields)
    for i in range(3):
        f_h = Inches(0.18)
        rect(s, col_x + Inches(1.9), y_fade + Inches(0.1) * i,
             Inches(2.9), Inches(0.1),
             fill=DEEP if i == 0 else (GHOST if i == 1 else SURFACE),
             stroke=None)

    # "and 65 more"
    text(s, col_x, y_fade + Inches(0.5), col_w, Inches(0.3),
         "\u2026 and 65 more.",
         font=MONO, size=10, italic=True, color=TERRACOTTA, bold=True,
         align=PP_ALIGN.CENTER)


# ============================================================
# SCENE 4 — THE IDEA
# ============================================================
def slide_the_idea():
    s = new_slide(bg=CANVAS)
    chrome(s, 3, "The idea")

    # kicker
    text(s, Inches(0.6), Inches(1.2), Inches(10), Inches(0.25),
         "A QUIET QUESTION",
         font=MONO, size=9, color=MUTED, bold=True)

    # short headline
    text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.8),
         "What if she could just\u2026 say it?",
         font=SERIF, size=40, italic=True, color=INK_SOFT)

    # centrepiece: a terracotta-bordered cream card, containing a giant
    # quote mark and the sentence that changes everything
    card_x = Inches(1.8)
    card_y = Inches(3.1)
    card_w = Inches(9.73)
    card_h = Inches(3.2)
    rect(s, card_x, card_y, card_w, card_h,
         fill=PAPER, stroke=TERRACOTTA, stroke_w=2.0)

    # giant opening quote mark
    text(s, card_x + Inches(0.2), card_y - Inches(0.3),
         Inches(2), Inches(2.5),
         "\u201C", font=SERIF, size=200, color=TERRACOTTA, line_spacing=1.0)

    # the sentence
    multi(s, card_x + Inches(1.6), card_y + Inches(0.65),
          card_w - Inches(2.0), card_h - Inches(1.0),
          ["3-bed ranch in OldTown,",
           "1,400 sqft, good shape."],
          font=SERIF, size=44, italic=True, color=INK, line_spacing=1.2)

    # decorative closing quote mark
    text(s, card_x + card_w - Inches(1.2), card_y + card_h - Inches(1.8),
         Inches(1.5), Inches(2.0),
         "\u201D", font=SERIF, size=120, color=TERRA_SOFT, line_spacing=1.0)

    # caption under card
    text(s, Inches(0.6), Inches(6.45), Inches(12.13), Inches(0.35),
         "One sentence. That\u2019s it. No fields, no dropdowns, no glossary.",
         font=SERIF, size=16, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SCENE 5 — THE TRANSFORMATION
# ============================================================
def slide_transformation():
    s = new_slide()
    chrome(s, 4, "The transformation")

    # kicker
    text(s, Inches(0.6), Inches(1.2), Inches(10), Inches(0.25),
         "WHAT HAPPENS NEXT",
         font=MONO, size=9, color=MUTED, bold=True)

    # headline
    multi(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(1.4),
          ["Three seconds later,"],
          font=SERIF, size=44, color=INK, line_spacing=1.0)
    multi(s, Inches(0.6), Inches(2.35), Inches(12.2), Inches(1.4),
          ["a number with a reason."],
          font=SERIF, size=44, italic=True, color=TERRACOTTA, line_spacing=1.0)

    # --- three-column layout: sentence | arrow | price ---
    row_y = Inches(4.15)
    row_h = Inches(2.4)

    # LEFT: the sentence card
    s_x = Inches(0.6)
    s_w = Inches(4.8)
    rect(s, s_x, row_y, s_w, row_h, fill=PAPER, stroke=HAIRLINE, stroke_w=1.0)
    text(s, s_x + Inches(0.3), row_y + Inches(0.3), s_w - Inches(0.6), Inches(0.3),
         "INPUT", font=MONO, size=9, color=MUTED, bold=True)
    hairline(s, s_x + Inches(0.3), row_y + Inches(0.6),
             s_x + Inches(1.1), row_y + Inches(0.6), color=TERRACOTTA, w=1.25)
    multi(s, s_x + Inches(0.3), row_y + Inches(0.85),
          s_w - Inches(0.6), row_h - Inches(1.1),
          ["\u201C3-bed ranch",
           "in OldTown,",
           "1,400 sqft,",
           "good shape.\u201D"],
          font=SERIF, size=16, italic=True, color=INK, line_spacing=1.3)

    # CENTER: the arrow
    a_x = Inches(5.6)
    a_w = Inches(2.1)
    arrow_y = row_y + row_h/2 - Inches(0.35)
    right_arrow(s, a_x, arrow_y, a_w, Inches(0.7), fill=TERRACOTTA)

    # time label above the arrow
    text(s, a_x - Inches(0.2), row_y + Inches(0.4),
         a_w + Inches(0.4), Inches(0.3),
         "\u2248 3 SECONDS", font=MONO, size=10, color=TERRACOTTA, bold=True,
         align=PP_ALIGN.CENTER)

    # pipeline label below the arrow
    multi(s, a_x - Inches(0.3), row_y + row_h - Inches(0.7),
          a_w + Inches(0.6), Inches(0.6),
          ["understand  \u00b7  predict", "\u00b7  explain"],
          font=MONO, size=9, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.4)

    # RIGHT: the price card
    p_x = Inches(7.9)
    p_w = Inches(4.83)
    rect(s, p_x, row_y, p_w, row_h, fill=PAPER, stroke=TERRACOTTA, stroke_w=1.5)
    text(s, p_x + Inches(0.3), row_y + Inches(0.3), p_w - Inches(0.6), Inches(0.3),
         "OUTPUT", font=MONO, size=9, color=TERRACOTTA, bold=True)
    hairline(s, p_x + Inches(0.3), row_y + Inches(0.6),
             p_x + Inches(1.2), row_y + Inches(0.6), color=TERRACOTTA, w=1.25)

    # big price
    text(s, p_x + Inches(0.3), row_y + Inches(0.75),
         Inches(0.8), Inches(1.3),
         "$", font=MONO, size=66, color=TERRACOTTA)
    text(s, p_x + Inches(0.95), row_y + Inches(0.75),
         p_w - Inches(1.2), Inches(1.3),
         "167,887", font=MONO, size=66, color=INK)

    # reasoning snippet
    text(s, p_x + Inches(0.3), row_y + Inches(1.95),
         p_w - Inches(0.6), Inches(0.35),
         "\u201CBelow the neighbourhood median \u2014",
         font=SERIF, size=12, italic=True, color=INK_SOFT)
    text(s, p_x + Inches(0.3), row_y + Inches(2.2),
         p_w - Inches(0.6), Inches(0.35),
         "smaller lot, older condition.\u201D",
         font=SERIF, size=12, italic=True, color=INK_SOFT)


# ============================================================
# SCENE 6 — UNDER THE HOOD (three helpers, non-technical)
# ============================================================
def slide_under_the_hood():
    s = new_slide()
    chrome(s, 5, "Under the hood")

    # kicker
    text(s, Inches(0.6), Inches(1.2), Inches(10), Inches(0.25),
         "HOW THE MAGIC ACTUALLY WORKS",
         font=MONO, size=9, color=MUTED, bold=True)

    # headline
    multi(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(1.3),
          ["Three helpers, working in sequence."],
          font=SERIF, size=36, color=INK, line_spacing=1.0)

    # narrator line
    text(s, Inches(0.6), Inches(2.7), Inches(12.13), Inches(0.3),
         "They don\u2019t sound technical \u2014 on purpose. You can describe them at a dinner party.",
         font=SERIF, size=14, italic=True, color=MUTED)

    # ---- three columns ----
    col_y = Inches(3.5)
    col_w = Inches(4.04)
    gap   = Inches(0.04)
    col_x0 = Inches(0.6)

    helpers = [
        {
            "num":   "01",
            "name":  "The Listener",
            "role":  "reads what she said",
            "body":  ["Turns her sentence into a list",
                      "of facts the computer can use.",
                      "Knows \u201Cranch\u201D means 1-storey."],
            "tech":  "GPT-4o-mini  +  Pydantic",
            "color": TERRACOTTA,
        },
        {
            "num":   "02",
            "name":  "The Expert",
            "role":  "has seen it before",
            "body":  ["Has looked at 2,051 real home",
                      "sales and learned the patterns.",
                      "Gives the price."],
            "tech":  "Random Forest  \u00b7  sklearn",
            "color": TEAL,
        },
        {
            "num":   "03",
            "name":  "The Narrator",
            "role":  "explains why",
            "body":  ["Takes the price + the facts",
                      "and writes a short reason",
                      "in plain English."],
            "tech":  "GPT-4o-mini  +  training stats",
            "color": BRASS,
        },
    ]

    for i, h in enumerate(helpers):
        x = col_x0 + (col_w + gap) * i

        # card background
        rect(s, x, col_y, col_w, Inches(3.1),
             fill=PAPER, stroke=HAIRLINE, stroke_w=0.75)

        # helper avatar (person silhouette in a circle)
        bg_r = Inches(1.25)
        bg_cx = x + col_w/2
        bg_cy = col_y + Inches(0.9)
        oval(s, bg_cx - bg_r/2, bg_cy - bg_r/2, bg_r, bg_r,
             fill=SURFACE)
        person(s, bg_cx, bg_cy, Inches(1.05), color=h["color"])

        # number
        text(s, x, col_y + Inches(1.75), col_w, Inches(0.3),
             h["num"], font=MONO, size=10, color=h["color"], bold=True,
             align=PP_ALIGN.CENTER)

        # name
        text(s, x, col_y + Inches(2.0), col_w, Inches(0.4),
             h["name"], font=SERIF, size=22, italic=True, color=INK,
             align=PP_ALIGN.CENTER)

        # role (small italic)
        text(s, x, col_y + Inches(2.45), col_w, Inches(0.3),
             h["role"], font=SERIF, size=13, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER)

        # description below
        multi(s, x + Inches(0.3), col_y + Inches(3.25),
              col_w - Inches(0.6), Inches(1.1),
              h["body"],
              font=SANS, size=12, color=INK_SOFT,
              align=PP_ALIGN.CENTER, line_spacing=1.4)

        # tech footnote
        text(s, x + Inches(0.2), col_y + Inches(3.0) - Inches(0.0),
             col_w - Inches(0.4), Inches(0.25),
             h["tech"], font=MONO, size=8, color=h["color"],
             align=PP_ALIGN.CENTER)

        # connecting arrow to next helper
        if i < 2:
            ax = x + col_w + gap/2 - Inches(0.12)
            ay = col_y + Inches(0.85)
            right_arrow(s, ax, ay, Inches(0.24), Inches(0.18),
                        fill=HAIRLINE)


# ============================================================
# SCENE 7 — THE PROOF
# ============================================================
def slide_the_proof():
    s = new_slide()
    chrome(s, 6, "The proof")

    # kicker
    text(s, Inches(0.6), Inches(1.2), Inches(10), Inches(0.25),
         "WHY LAYLA CAN TRUST IT",
         font=MONO, size=9, color=MUTED, bold=True)

    # headline
    multi(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(1.3),
          ["Not a party trick."],
          font=SERIF, size=44, color=INK, line_spacing=1.0)
    text(s, Inches(0.6), Inches(2.7), Inches(12.13), Inches(0.35),
         "Three numbers that earn her confidence before the first click.",
         font=SERIF, size=14, italic=True, color=MUTED)

    # ---- three stat blocks ----
    row_y = Inches(3.5)
    row_h = Inches(3.1)
    gap   = Inches(0.3)
    cw    = (Inches(12.13) - gap * 2) / 3
    cx0   = Inches(0.6)

    stats = [
        {
            "icon":   "houses",
            "hero":   "2,051",
            "unit":   "real homes",
            "body":   "Every estimate is rooted in actual sales from Ames, Iowa.",
            "tech":   "training set, 70/15/15 split",
            "color":  TEAL,
        },
        {
            "icon":   "target",
            "hero":   "~$25K",
            "unit":   "typical miss",
            "body":   "About one kitchen renovation \u2014 on a dataset that ranges from $12K to $755K.",
            "tech":   "RMSE on held-out test  \u00b7  R\u00b2 0.896",
            "color":  TERRACOTTA,
        },
        {
            "icon":   "chat",
            "hero":   "98%",
            "unit":   "understood",
            "body":   "Nearly every normal sentence becomes the right structured facts.",
            "tech":   "prompt v2  \u00b7  labeled benchmark",
            "color":  BRASS,
        },
    ]

    for i, st in enumerate(stats):
        x = cx0 + (cw + gap) * i
        # card
        rect(s, x, row_y, cw, row_h, fill=PAPER, stroke=HAIRLINE, stroke_w=0.75)

        # icon area
        icon_cx = x + cw/2
        icon_cy = row_y + Inches(0.75)
        if st["icon"] == "houses":
            house_stack_icon(s, icon_cx, icon_cy, Inches(2.5), Inches(0.9),
                             color=st["color"])
        elif st["icon"] == "target":
            target_icon(s, icon_cx, icon_cy, Inches(0.95), color=st["color"])
        elif st["icon"] == "chat":
            chat_icon(s, icon_cx, icon_cy, Inches(1.3), Inches(0.95),
                      color=st["color"])

        # hero number
        text(s, x, row_y + Inches(1.4), cw, Inches(0.9),
             st["hero"], font=SERIF, size=60, bold=True, color=INK,
             align=PP_ALIGN.CENTER)
        # unit
        text(s, x, row_y + Inches(2.3), cw, Inches(0.3),
             st["unit"], font=SERIF, size=16, italic=True, color=st["color"],
             align=PP_ALIGN.CENTER)
        # hairline
        hairline(s, x + Inches(1.3), row_y + Inches(2.68),
                 x + cw - Inches(1.3), row_y + Inches(2.68),
                 color=HAIRLINE, w=0.5)
        # body
        multi(s, x + Inches(0.3), row_y + Inches(2.78),
              cw - Inches(0.6), Inches(0.7),
              [st["body"]],
              font=SANS, size=11, color=INK_SOFT,
              align=PP_ALIGN.CENTER, line_spacing=1.35)

    # tech footnotes row (under all three cards)
    tech_y = row_y + row_h + Inches(0.08)
    for i, st in enumerate(stats):
        x = cx0 + (cw + gap) * i
        text(s, x, tech_y, cw, Inches(0.22),
             st["tech"], font=MONO, size=8, color=MUTED,
             align=PP_ALIGN.CENTER)


# ============================================================
# SCENE 8 — THE DEMO
# ============================================================
def slide_demo():
    s = new_slide()
    chrome(s, 7, "The demo")

    # kicker
    text(s, Inches(0.6), Inches(1.2), Inches(10), Inches(0.25),
         "LET\u2019S WATCH LAYLA USE IT",
         font=MONO, size=9, color=MUTED, bold=True)

    # headline
    multi(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.2),
          ["Now", "she just", "says it."],
          font=SERIF, size=50, color=INK, line_spacing=1.02)

    hairline(s, Inches(0.6), Inches(4.4), Inches(6.1), Inches(4.4),
             color=TERRACOTTA, w=1.25)

    # speaker cues (small)
    text(s, Inches(0.6), Inches(4.55), Inches(5.5), Inches(0.25),
         "SPEAKER CUES", font=MONO, size=9, color=MUTED, bold=True)

    cues = [
        ("01", "Type the OldTown example."),
        ("02", "Let the three stages light up."),
        ("03", "Read the price aloud."),
        ("04", "Read one line of the reasoning."),
        ("05", "If time allows, try luxury NridgHt."),
    ]
    y = Inches(4.9)
    for num, body in cues:
        text(s, Inches(0.6), y, Inches(0.4), Inches(0.25),
             num, font=MONO, size=10, color=TERRACOTTA, bold=True)
        text(s, Inches(1.0), y, Inches(5.1), Inches(0.25),
             body, font=MONO, size=10, color=INK)
        y += Inches(0.25)

    # ---- RIGHT: the product mockup ----
    fx = Inches(6.7)
    fy = Inches(1.3)
    fw = Inches(6.15)
    fh = Inches(5.55)

    # browser frame
    rect(s, fx, fy, fw, fh, fill=PAPER, stroke=HAIRLINE, stroke_w=1.0)

    # top chrome
    rect(s, fx, fy, fw, Inches(0.4), fill=SURFACE, stroke=HAIRLINE, stroke_w=0.5)
    hairline(s, fx, fy + Inches(0.4), fx + fw, fy + Inches(0.4), color=HAIRLINE)

    d = Inches(0.12)
    oval(s, fx + Inches(0.2),  fy + Inches(0.13), d, d, fill=TERRACOTTA)
    oval(s, fx + Inches(0.4),  fy + Inches(0.13), d, d, fill=BRASS)
    oval(s, fx + Inches(0.6),  fy + Inches(0.13), d, d, fill=TEAL)

    url_w = Inches(3.0)
    url_x = fx + (fw - url_w) / 2
    rrect(s, url_x, fy + Inches(0.1), url_w, Inches(0.22),
          fill=PAPER, stroke=HAIRLINE, stroke_w=0.5, radius=0.2)
    text(s, url_x, fy + Inches(0.11), url_w, Inches(0.2),
         "aire.local  /  predict", font=MONO, size=9, color=MUTED,
         align=PP_ALIGN.CENTER)

    # stage strip
    st_y = fy + Inches(0.6)
    st_h = Inches(0.55)
    st_pad = Inches(0.35)
    st_w = fw - st_pad * 2
    rect(s, fx + st_pad, st_y, st_w, st_h, fill=SURFACE)
    hairline(s, fx + st_pad, st_y + st_h, fx + st_pad + st_w, st_y + st_h,
             color=HAIRLINE, w=0.5)
    stage_w = st_w / 3
    stage_labels = [("01", "LISTEN"), ("02", "PREDICT"), ("03", "EXPLAIN")]
    for i, (n, lab) in enumerate(stage_labels):
        sx = fx + st_pad + stage_w * i
        oval(s, sx + Inches(0.18), st_y + Inches(0.18),
             Inches(0.18), Inches(0.18), fill=TEAL)
        text(s, sx + Inches(0.18), st_y + Inches(0.15),
             Inches(0.18), Inches(0.22),
             "\u2713", font=SANS, size=10, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER)
        text(s, sx + Inches(0.48), st_y + Inches(0.18),
             stage_w - Inches(0.5), Inches(0.22),
             f"{n}  {lab}", font=MONO, size=9, color=INK, bold=True)
        if i < 2:
            hairline(s, sx + stage_w, st_y + Inches(0.1),
                     sx + stage_w, st_y + st_h - Inches(0.1),
                     color=HAIRLINE, w=0.5)

    # content
    cx = fx + Inches(0.45)
    cy = fy + Inches(1.35)
    cw = fw - Inches(0.9)

    text(s, cx, cy, cw, Inches(0.25),
         "04   ESTIMATE", font=MONO, size=10, color=TERRACOTTA, bold=True)
    hairline(s, cx, cy + Inches(0.3), cx + Inches(0.9), cy + Inches(0.3),
             color=TERRACOTTA, w=1.25)

    # big price
    text(s, cx - Inches(0.05), cy + Inches(0.4), Inches(0.9), Inches(1.7),
         "$", font=MONO, size=78, color=TERRACOTTA)
    text(s, cx + Inches(0.7), cy + Inches(0.4), cw - Inches(0.9), Inches(1.7),
         "167,887", font=MONO, size=78, color=INK)
    text(s, cx, cy + Inches(1.9), cw, Inches(0.25),
         "USD  \u00b7  point estimate",
         font=MONO, size=9, color=MUTED)

    # reasoning
    hairline(s, cx, cy + Inches(2.3), cx + cw, cy + Inches(2.3), color=HAIRLINE)
    text(s, cx, cy + Inches(2.4), cw, Inches(0.25),
         "REASONING", font=MONO, size=9, color=MUTED, bold=True)
    multi(s, cx, cy + Inches(2.65), cw, Inches(1.0),
          ["\u201CBelow the $160,000 neighbourhood median.",
           "A smaller lot and an older overall condition",
           "pull the estimate down.\u201D"],
          font=SERIF, size=13, italic=True, color=INK, line_spacing=1.35)

    # feature chips
    hairline(s, cx, cy + Inches(4.0), cx + cw, cy + Inches(4.0), color=HAIRLINE)
    text(s, cx, cy + Inches(4.1), cw, Inches(0.25),
         "UNDERSTOOD  \u00b7  4 of 12", font=MONO, size=9, color=MUTED, bold=True)

    chips = [
        ("BedroomAbvGr", "3"),
        ("Neighborhood", "OldTown"),
        ("GrLivArea",    "1,400"),
        ("OverallQual",  "5*"),
    ]
    cg = Inches(0.08)
    chw = (cw - cg * 3) / 4
    chy = cy + Inches(4.35)
    chh = Inches(0.55)
    for i, (k, v) in enumerate(chips):
        cx_i = cx + (chw + cg) * i
        rrect(s, cx_i, chy, chw, chh, fill=SURFACE, stroke=HAIRLINE,
              stroke_w=0.5, radius=0.15)
        text(s, cx_i + Inches(0.1), chy + Inches(0.05),
             chw - Inches(0.2), Inches(0.22),
             k, font=MONO, size=8, color=MUTED)
        text(s, cx_i + Inches(0.1), chy + Inches(0.25),
             chw - Inches(0.2), Inches(0.3),
             v, font=SERIF, size=14, italic=True, color=INK)

    text(s, cx, chy + chh + Inches(0.07), cw, Inches(0.22),
         "*  imputed from training-set median",
         font=MONO, size=8, color=MUTED)


# ============================================================
# SCENE 9 — THE VISION (CLOSING)
# ============================================================
def slide_vision():
    s = new_slide()
    chrome(s, 8, "The vision")

    # kicker
    text(s, Inches(0.6), Inches(1.2), Inches(10), Inches(0.25),
         "WHERE THIS GOES",
         font=MONO, size=9, color=MUTED, bold=True)

    # small section header
    text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.55),
         "Shipped, end-to-end.", font=SERIF, size=28, italic=True, color=INK_SOFT)

    # three checked lines
    hairline(s, Inches(0.6), Inches(2.4), Inches(12.73), Inches(2.4),
             color=HAIRLINE)

    items = [
        ("A language model listens.  A real model prices.  A second language model explains.",
         "app/chain/  \u00b7  app/ml/  \u00b7  app/main.py"),
        ("The winning prompt was chosen against labelled data.",
         "v2 accuracy 0.983  \u00b7  v1 accuracy 0.617"),
        ("One command deploys it \u2014 backend, frontend, model, all boxed up.",
         "docker compose up"),
    ]
    y = Inches(2.6)
    for body, detail in items:
        oval(s, Inches(0.6), y + Inches(0.08),
             Inches(0.22), Inches(0.22), fill=TEAL)
        text(s, Inches(0.6), y + Inches(0.04),
             Inches(0.22), Inches(0.25),
             "\u2713", font=SANS, size=11, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER)
        text(s, Inches(1.0), y, Inches(11.5), Inches(0.35),
             body, font=SANS, size=15, color=INK)
        text(s, Inches(1.0), y + Inches(0.38), Inches(11.5), Inches(0.25),
             detail, font=MONO, size=9, color=MUTED)
        y += Inches(0.7)

    # NEXT strip
    hairline(s, Inches(0.6), y + Inches(0.05), Inches(12.73), y + Inches(0.05),
             color=HAIRLINE)
    text(s, Inches(0.6), y + Inches(0.22), Inches(1.3), Inches(0.3),
         "NEXT", font=MONO, size=10, color=TERRACOTTA, bold=True)
    text(s, Inches(1.6), y + Inches(0.2), Inches(11.13), Inches(0.35),
         "Comparable listings.  Confidence around every price.",
         font=SERIF, size=15, italic=True, color=INK)

    # --- closing line — biggest type in the deck ---
    hairline(s, Inches(4.5), Inches(6.05), Inches(8.83), Inches(6.05),
             color=TERRACOTTA, w=2.0)
    multi(s, Inches(0.6), Inches(6.2), Inches(12.13), Inches(0.9),
          ["A price with a reason",
           "is no longer a luxury."],
          font=SERIF, size=40, color=INK, italic=True,
          line_spacing=1.02, align=PP_ALIGN.CENTER)


# ============================================================
# BUILD
# ============================================================
slide_cover()
slide_meet_sarah()
slide_the_wall()
slide_the_idea()
slide_transformation()
slide_under_the_hood()
slide_the_proof()
slide_demo()
slide_vision()

out = "AI_Real_Estate_Agent.pptx"
prs.save(out)
print(f"wrote {out}  ·  {len(prs.slides)} scenes  ·  fade transitions applied")
